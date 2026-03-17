"""
Tests for export-pptx.py

Run: python -m pytest tests/test_pptx.py -v
     python tests/run_tests.py
"""

import subprocess
import sys
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

FIXTURES = Path(__file__).parent / "fixtures"
SCRIPT = Path(__file__).parent.parent / "scripts" / "export-pptx.py"


def run_export(args: list, cwd=None) -> subprocess.CompletedProcess:
    return subprocess.run(
        [sys.executable, str(SCRIPT)] + args,
        capture_output=True, text=True, cwd=cwd
    )


# ─── Happy path ──────────────────────────────────────────────────────────────

class TestPptxExport:

    def test_slide_count_matches(self, tmp_path):
        """3 .slide elements → PPTX should have exactly 3 slides."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        assert out.exists(), "Output PPTX not created"
        prs = Presentation(str(out))
        assert len(prs.slides) == 3

    def test_output_file_has_content(self, tmp_path):
        """Output PPTX should be larger than 10 KB (proves real screenshots)."""
        out = tmp_path / "out.pptx"
        run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        assert out.stat().st_size > 10_000, "PPTX is suspiciously small"

    def test_default_output_path(self, tmp_path):
        """No output arg → PPTX is saved alongside the HTML with same stem."""
        import shutil
        html = tmp_path / "my_deck.html"
        shutil.copy(FIXTURES / "simple_slides.html", html)
        r = run_export([str(html)])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        pptx = tmp_path / "my_deck.pptx"
        assert pptx.exists(), f"Expected {pptx} to be created"

    def test_default_aspect_ratio_matches_viewport(self, tmp_path):
        """Default 1440×900 viewport → slide ratio should match 1440/900 = 1.6 (8:5)."""
        out = tmp_path / "out.pptx"
        run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        prs = Presentation(str(out))
        ratio = prs.slide_width / prs.slide_height
        expected = 1440 / 900  # 1.6 (8:5), not 16:9
        assert abs(ratio - expected) < 0.01, f"Expected {expected:.4f}, got {ratio:.4f}"

    def test_custom_dimensions(self, tmp_path):
        """--width 1920 --height 1080 → slide ratio should be 16:9."""
        out = tmp_path / "out.pptx"
        r = run_export([
            str(FIXTURES / "simple_slides.html"), str(out),
            "--width", "1920", "--height", "1080"
        ])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        prs = Presentation(str(out))
        assert len(prs.slides) == 3
        ratio = prs.slide_width / prs.slide_height
        assert abs(ratio - 16 / 9) < 0.01

    def test_custom_4x3_ratio(self, tmp_path):
        """--width 1024 --height 768 → slide ratio should be 4:3."""
        out = tmp_path / "out.pptx"
        run_export([
            str(FIXTURES / "simple_slides.html"), str(out),
            "--width", "1024", "--height", "768"
        ])
        prs = Presentation(str(out))
        ratio = prs.slide_width / prs.slide_height
        assert abs(ratio - 4 / 3) < 0.01, f"Expected 4:3, got {ratio:.4f}"

    def test_stdout_reports_slide_count(self, tmp_path):
        """Script stdout should mention the number of slides found."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        assert "3" in r.stdout, f"Expected slide count in stdout: {r.stdout}"


# ─── Error cases ─────────────────────────────────────────────────────────────

class TestPptxErrorCases:

    def test_no_slide_elements_exits_gracefully(self, tmp_path):
        """HTML with no .slide elements → script exits 0 with 'Nothing to export'."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "no_slides.html"), str(out)])
        # Should exit cleanly (not crash), output file should not exist
        assert r.returncode == 0
        assert not out.exists(), "No PPTX should be created when there are no slides"
        assert "Nothing" in r.stdout or "0 slides" in r.stdout.lower() or "No .slide" in r.stdout

    def test_missing_file_exits_nonzero(self, tmp_path):
        """Nonexistent input file → non-zero exit code."""
        out = tmp_path / "out.pptx"
        r = run_export([str(tmp_path / "ghost.html"), str(out)])
        assert r.returncode != 0, "Expected non-zero exit for missing file"
